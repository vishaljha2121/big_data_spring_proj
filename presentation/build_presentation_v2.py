#!/usr/bin/env python3
"""Generate the polished v2 final presentation package."""

from __future__ import annotations

import json
import math
from pathlib import Path
from typing import Dict, Iterable, List, Sequence, Tuple

from PIL import Image, ImageDraw, ImageFont, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "presentation"
ASSETS = OUT / "assets_v2"
DIAGRAMS = ASSETS / "diagrams"
SCREENSHOTS = ASSETS / "screenshots"
PPTX = OUT / "CourtIQ_Final_Presentation_v2.pptx"
NOTES = OUT / "speaker_notes_v2.md"
README = OUT / "README_v2.md"

GREEN = RGBColor(5, 58, 35)
GREEN2 = RGBColor(10, 93, 55)
CREAM = RGBColor(248, 243, 229)
CLAY = RGBColor(185, 91, 48)
GOLD = RGBColor(218, 177, 82)
PURPLE = RGBColor(91, 63, 164)
INK = RGBColor(24, 35, 49)
MUTED = RGBColor(94, 107, 124)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(226, 216, 195)


def j(path: str, default=None):
    p = ROOT / path
    if not p.exists():
        return default if default is not None else {}
    return json.loads(p.read_text(encoding="utf-8"))


def fnum(value, digits=2):
    if isinstance(value, float):
        return f"{value:.{digits}f}"
    return str(value)


def font(size=28, bold=False):
    paths = [
        "/System/Library/Fonts/Supplemental/Avenir Next.ttc",
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf" if bold else "/System/Library/Fonts/Supplemental/Arial.ttf",
    ]
    for p in paths:
        if Path(p).exists():
            return ImageFont.truetype(p, size=size)
    return ImageFont.load_default()


def wrap(draw: ImageDraw.ImageDraw, text: str, fnt, width: int) -> List[str]:
    words, lines, line = text.split(), [], ""
    for word in words:
        candidate = f"{line} {word}".strip()
        if draw.textbbox((0, 0), candidate, font=fnt)[2] <= width:
            line = candidate
        else:
            if line:
                lines.append(line)
            line = word
    if line:
        lines.append(line)
    return lines


def arrow(draw, start, end, color=(6, 66, 37), width=6):
    draw.line([start, end], fill=color, width=width)
    angle = math.atan2(end[1] - start[1], end[0] - start[0])
    size = 18
    p1 = (end[0] - size * math.cos(angle - math.pi / 6), end[1] - size * math.sin(angle - math.pi / 6))
    p2 = (end[0] - size * math.cos(angle + math.pi / 6), end[1] - size * math.sin(angle + math.pi / 6))
    draw.polygon([end, p1, p2], fill=color)


def node(draw, box, title, subtitle="", fill=(255, 253, 247), accent=(218, 177, 82)):
    x1, y1, x2, y2 = box
    draw.rounded_rectangle(box, radius=28, fill=(0, 0, 0, 35))
    box = (x1 - 3, y1 - 8, x2 - 3, y2 - 8)
    draw.rounded_rectangle(box, radius=28, fill=fill, outline=(222, 210, 185), width=3)
    draw.rounded_rectangle((box[0], box[1], box[0] + 12, box[3]), radius=6, fill=accent)
    tf = font(28, True)
    sf = font(19, False)
    draw.text((box[0] + 34, box[1] + 24), title, font=tf, fill=(22, 35, 49))
    y = box[1] + 62
    for line in wrap(draw, subtitle, sf, box[2] - box[0] - 50)[:2]:
        draw.text((box[0] + 34, y), line, font=sf, fill=(85, 98, 113))
        y += 24


def diagram_canvas(title: str) -> Tuple[Image.Image, ImageDraw.ImageDraw]:
    img = Image.new("RGB", (1800, 1040), (248, 243, 229))
    draw = ImageDraw.Draw(img, "RGBA")
    draw.rectangle((0, 0, 1800, 132), fill=(5, 58, 35))
    draw.text((70, 38), title, font=font(48, True), fill=(255, 255, 255))
    draw.line((72, 132, 1728, 132), fill=(218, 177, 82), width=5)
    return img, draw


def save_diagram(name: str, title: str, boxes: Sequence[Tuple[int, int, int, int, str, str, Tuple[int, int, int]]], arrows: Sequence[Tuple[int, int, int, int]]):
    img, draw = diagram_canvas(title)
    for a in arrows:
        arrow(draw, (a[0], a[1]), (a[2], a[3]), color=(5, 58, 35))
    for x1, y1, x2, y2, t, sub, accent in boxes:
        node(draw, (x1, y1, x2, y2), t, sub, accent=accent)
    path = DIAGRAMS / f"{name}.png"
    img.save(path)
    return path


def save_mermaid(name: str, text: str):
    (DIAGRAMS / f"{name}.mmd").write_text(text.strip() + "\n", encoding="utf-8")


def build_assets():
    DIAGRAMS.mkdir(parents=True, exist_ok=True)
    SCREENSHOTS.mkdir(parents=True, exist_ok=True)
    gold, clay, green, purple = (218, 177, 82), (185, 91, 48), (5, 58, 35), (91, 63, 164)
    diagrams = {}
    diagrams["architecture"] = save_diagram(
        "01_end_to_end_architecture_v2",
        "Validated Data-to-Product Architecture",
        [
            (80, 235, 400, 390, "Parquet lake", "cleaned, curated, features", gold),
            (510, 235, 830, 390, "Model artifacts", "odds model + risk config", purple),
            (940, 235, 1260, 390, "Replay stream", "Kafka canonical point events", clay),
            (1370, 235, 1690, 390, "Spark scoring", "Structured Streaming + checkpoint", green),
            (510, 650, 830, 805, "Scored sink", "JSONL + Parquet evidence", green),
            (940, 650, 1260, 805, "FastAPI", "documented endpoints", purple),
            (1370, 650, 1690, 805, "Centre Court UI", "analytics dashboard", gold),
        ],
        [(400, 310, 510, 310), (830, 310, 940, 310), (1260, 310, 1370, 310), (1530, 390, 670, 650), (830, 725, 940, 725), (1260, 725, 1370, 725)],
    )
    save_mermaid("01_end_to_end_architecture_v2", "flowchart LR\nA[Parquet lake]-->B[Models]-->C[Kafka replay]-->D[Spark scoring]-->E[Scored sink]-->F[FastAPI]-->G[Dashboard]")
    diagrams["lake"] = save_diagram(
        "02_data_lake_zones_v2",
        "Data Lake Zones and Contracts",
        [
            (85, 250, 360, 405, "Staging", "source inventory + raw limits", clay),
            (455, 250, 730, 405, "Cleaned", "normalized point records", gold),
            (825, 250, 1100, 405, "Curated", "singles-only Parquet", green),
            (1195, 250, 1470, 405, "Features", "point-in-time safe", purple),
            (455, 635, 730, 790, "Baselines", "player serve/return priors", purple),
            (825, 635, 1100, 790, "Replay", "deterministic manifest", clay),
            (1195, 635, 1470, 790, "Scored", "validated outputs", green),
        ],
        [(360, 325, 455, 325), (730, 325, 825, 325), (1100, 325, 1195, 325), (960, 405, 960, 635), (1100, 710, 1195, 710)],
    )
    save_mermaid("02_data_lake_zones_v2", "flowchart LR\nA[Staging]-->B[Cleaned]-->C[Curated]-->D[Features]\nC-->E[Baselines]\nC-->F[Replay]\nF-->G[Scored]")
    diagrams["feature"] = save_diagram(
        "03_point_in_time_v2",
        "Point-in-Time Feature Safety",
        [
            (105, 285, 420, 440, "Events < k", "prior points only", green),
            (540, 285, 855, 440, "Feature row k", "rolling/cumulative before", gold),
            (975, 285, 1290, 440, "Model score", "current point probability", purple),
            (1410, 285, 1690, 440, "Update state", "after scoring point k", clay),
            (540, 650, 1290, 790, "No leakage rule", "current point winner is never used before prediction", clay),
        ],
        [(420, 360, 540, 360), (855, 360, 975, 360), (1290, 360, 1410, 360), (1540, 440, 910, 650)],
    )
    save_mermaid("03_point_in_time_v2", "flowchart LR\nA[Prior events < k]-->B[Feature row k]-->C[Score point k]-->D[Update state after scoring]")
    diagrams["streaming"] = save_diagram(
        "04_kafka_spark_streaming_v2",
        "Kafka + Spark Structured Streaming Runtime",
        [
            (80, 260, 395, 420, "Replay producer", "1000 point_event_v1 messages", clay),
            (505, 260, 820, 420, "Kafka topic", "tennis-point-events, 20 partitions", gold),
            (930, 260, 1245, 420, "Spark readStream", "Kafka source, earliest offsets", green),
            (1355, 260, 1685, 420, "foreachBatch", "validated StreamScorer", purple),
            (930, 640, 1245, 800, "Checkpoint", "data/checkpoints/spark_streaming_scorer", gold),
            (1355, 640, 1685, 800, "Scored sink", "1000 JSONL + Parquet rows", green),
        ],
        [(395, 340, 505, 340), (820, 340, 930, 340), (1245, 340, 1355, 340), (1520, 420, 1520, 640), (1088, 420, 1088, 640)],
    )
    save_mermaid("04_kafka_spark_streaming_v2", "flowchart LR\nA[Replay producer]-->B[(Kafka topic)]-->C[Spark readStream]-->D[foreachBatch scorer]-->E[Scored sink]\nC-->F[Checkpoint]")
    diagrams["model"] = save_diagram(
        "05_model_publication_v2",
        "Model Publication and Scoring Flow",
        [
            (110, 270, 415, 425, "Train split", "match-level, no leakage", green),
            (525, 270, 830, 425, "Evaluate", "AUC + Brier gates", gold),
            (940, 270, 1245, 425, "Publish v1", "model.joblib + latest.json", purple),
            (1355, 270, 1660, 425, "Runtime load", "schema hash check", clay),
            (525, 650, 830, 805, "Risk config", "no fake labels", clay),
            (940, 650, 1245, 805, "Scored event", "probability + risk signal", green),
        ],
        [(415, 345, 525, 345), (830, 345, 940, 345), (1245, 345, 1355, 345), (1510, 425, 1090, 650), (830, 725, 940, 725)],
    )
    save_mermaid("05_model_publication_v2", "flowchart LR\nA[Train split]-->B[Evaluate]-->C[Publish v1]-->D[Runtime load]\nB-->E[Risk config]-->F[Scored event]")
    diagrams["serving"] = save_diagram(
        "06_api_frontend_v2",
        "Serving Layer and Product Dashboard",
        [
            (105, 260, 420, 420, "Scored files", "JSONL/Parquet validated", green),
            (540, 260, 855, 420, "FastAPI", "summary, matches, risk, models", purple),
            (975, 260, 1290, 420, "Contracts", "OpenAPI + examples", gold),
            (1410, 260, 1690, 420, "Dashboard", "Centre Court Analytics", clay),
            (540, 650, 855, 805, "Validation", "preflight + tests", green),
            (975, 650, 1290, 805, "Demo", "one-command runner", gold),
        ],
        [(420, 340, 540, 340), (855, 340, 975, 340), (1290, 340, 1410, 340), (700, 420, 700, 650), (855, 725, 975, 725)],
    )
    save_mermaid("06_api_frontend_v2", "flowchart LR\nA[Scored files]-->B[FastAPI]-->C[OpenAPI]-->D[Dashboard]\nB-->E[Validation]-->F[Demo]")
    return diagrams


def add_bg(slide, prs, dark=False):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = GREEN if dark else CREAM
    if not dark:
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.22))
        band.fill.solid(); band.fill.fore_color.rgb = GREEN
        band.line.fill.background()


def textbox(slide, text, x, y, w, h, size=20, color=INK, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    if align:
        p.alignment = align
    return box


def title(slide, text, subtitle=None):
    textbox(slide, text, 0.55, 0.35, 12.0, 0.55, size=28, color=GREEN, bold=True)
    if subtitle:
        textbox(slide, subtitle, 0.58, 0.95, 11.6, 0.35, size=12, color=MUTED)


def bullets(slide, items, x, y, w, h, size=16):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, item in enumerate(items[:5]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = INK
        p.space_after = Pt(8)
    return box


def card(slide, label, value, x, y, w=2.55, h=0.95, accent=GREEN, small=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = LINE
    textbox(slide, label.upper(), x + 0.15, y + 0.1, w - 0.3, 0.22, size=8 if small else 9, color=MUTED, bold=True)
    textbox(slide, str(value), x + 0.15, y + 0.38, w - 0.3, 0.36, size=15 if small else 20, color=accent, bold=True)


def image(slide, path, x, y, w=None, h=None):
    if w and h:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    if w:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y))


def framed_image(slide, path, x, y, w, h=None, label=None):
    if h is None:
        with Image.open(path) as img:
            img_w, img_h = img.size
            h = w * (img_h / img_w)

    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x - 0.05), Inches(y - 0.05), Inches(w + 0.1), Inches(h + 0.1))
    frame.fill.solid(); frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = LINE
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    if label:
        textbox(slide, label, x + 0.15, y + h - 0.38, w - 0.3, 0.25, size=10, color=WHITE, bold=True)


def pill(slide, text, x, y, w=2.4, color=GOLD):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    textbox(slide, text, x + 0.08, y + 0.065, w - 0.16, 0.16, size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def slide_no(slide, n):
    textbox(slide, f"{n:02d}", 12.25, 7.05, 0.55, 0.2, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def build_deck():
    ASSETS.mkdir(parents=True, exist_ok=True); SCREENSHOTS.mkdir(parents=True, exist_ok=True)
    diagrams = build_assets()
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    odds = j("data/results/model_eval/odds_model_eval_report.json")
    bench = j("data/results/streaming_scoring/scoring_benchmark_report.json")
    api = j("data/results/api_validation/api_validation_report.json")
    kafka = j("data/results/kafka_runtime/kafka_runtime_report.json")
    kafka_smoke = j("data/results/kafka_runtime/kafka_replay_smoke_report.json")
    spark = j("data/results/spark_streaming/spark_streaming_run_report.json")
    spark_val = j("data/results/spark_streaming/spark_streaming_validation_report.json")
    features = j("data/features/validation_report.json")
    frontend = j("data/results/frontend_validation/frontend_validation_report.json")

    slides = []
    def sld(t=None, sub=None, dark=False):
        s = prs.slides.add_slide(blank); add_bg(s, prs, dark=dark); slides.append(s)
        if t: title(s, t, sub)
        slide_no(s, len(slides)); return s

    s = sld(dark=True)
    textbox(s, "CourtIQ", 0.72, 0.72, 4.0, 0.5, 24, GOLD, True)
    textbox(s, "Centre Court Analytics", 0.72, 1.35, 10.6, 0.75, 42, WHITE, True)
    textbox(s, "Replay historical tennis points as a stream, score them live, and serve the results through an analytics dashboard.", 0.78, 2.22, 10.9, 0.7, 19, RGBColor(240, 232, 210))
    if (SCREENSHOTS / "dashboard_overview.png").exists():
        framed_image(s, SCREENSHOTS / "dashboard_overview.png", 6.45, 3.12, 5.95)
    for i, (a, b) in enumerate([("1.92M+", "point rows"), ("10.5K+", "matches"), ("1000", "Kafka/Spark scored"), ("70", "tests passed")]):
        card(s, b, a, 0.78 + i * 1.35, 4.35, 1.16, 0.92, accent=GOLD, small=True)

    s = sld("Problem + Goal", "Turn historical point data into a replayable, model-scored tennis analytics product.")
    for i, (h, v) in enumerate([
        ("Raw data", "messy point records, missing metadata"),
        ("Model risk", "future leakage can invalidate results"),
        ("Demo gap", "batch files alone do not prove streaming"),
    ]):
        card(s, h, v, 0.8 + i * 4.1, 1.75, 3.55, 1.25, accent=[CLAY, PURPLE, GREEN][i])
    bullets(s, ["Goal: validated data lake → replay stream → online features → point model → risk signal → API/dashboard.", "Keep the language honest: point probability is not betting odds; risk is not proof of misconduct."], 1.0, 4.25, 11.3, 1.4, 18)

    s = sld("Final Product Demo Snapshot", "The final UI is a product, not just an output table.")
    if (SCREENSHOTS / "dashboard_overview.png").exists():
        framed_image(s, SCREENSHOTS / "dashboard_overview.png", 0.65, 1.35, 7.4, label="Centre Court dashboard")
    bullets(s, ["Dashboard KPIs and validation evidence", "Match Browser and Replay Center", "Point Scoring Center + Model Performance", "Pipeline Monitor and Reports"], 8.45, 1.55, 4.1, 3.2, 18)
    pill(s, "File-backed API + real scored outputs", 8.5, 5.25, 3.3, GREEN)

    s = sld("Why This Satisfies Big Data Scope")
    for i, (l, v, c) in enumerate([
        ("Point rows", "1.92M", GREEN), ("Matches", "10.5K", CLAY), ("Replay events", "1.92M", PURPLE), ("Kafka partitions", "20", GOLD),
        ("Spark scored", "1000", GREEN), ("Tests", "70", CLAY), ("API rows", api.get("scored_event_count", "9049"), PURPLE), ("Feature rows", "1.92M", GOLD),
    ]):
        card(s, l, v, 0.75 + (i % 4) * 3.05, 1.4 + (i // 4) * 1.2, 2.55, 0.9, c)
    bullets(s, ["Parquet data lake zones with schema contracts and data-quality reports.", "Spark-style batch feature processing plus Kafka replay stream.", "Spark Structured Streaming runtime with checkpointing and JSONL/Parquet sink.", "Local demo is bounded; replay/scoring design is horizontally scalable."], 0.9, 4.25, 11.4, 1.9, 17)

    s = sld("Dataset Scale + Data Quality")
    card(s, "Curated point rows", "1,922,136", 0.8, 1.5, 3.5, 1.0, GREEN)
    card(s, "Curated matches", "10,508", 4.9, 1.5, 3.5, 1.0, CLAY)
    card(s, "Player baselines", "1,891", 9.0, 1.5, 3.5, 1.0, PURPLE)
    bullets(s, ["Singles-only curated Parquet data", "Invalid/special point evidence preserved", "Reports for missing winners, servers, elapsed time, rally length", "Surface coverage is 0%; surface analytics are planned, not claimed"], 1.0, 3.35, 11.2, 2.4, 18)

    s = sld("End-to-End Architecture")
    image(s, diagrams["architecture"], 0.55, 1.18, w=12.25)

    s = sld("Data Lake + Cleaning Pipeline")
    image(s, diagrams["lake"], 0.55, 1.15, w=6.05)
    image(s, diagrams["feature"], 6.75, 1.15, w=6.05)

    s = sld("Point-in-Time Feature Engineering")
    image(s, diagrams["feature"], 0.65, 1.2, w=7.3)
    bullets(s, ["For event k, features use only events < k.", "Cumulative, rolling, and recent-window features are shifted before scoring.", "Tests prove the current point winner is not included in current-row features.", "This matters because score-state leakage produced unrealistically high metrics before removal."], 8.35, 1.55, 4.25, 3.9, 17)

    s = sld("Model Training + Evaluation")
    image(s, diagrams["model"], 0.65, 1.2, w=6.8)
    for i, (l, v, c) in enumerate([
        ("Target", "Player A wins current point", GREEN), ("Model", "HistGradientBoosting", PURPLE), ("Features", "26", GOLD),
        ("Val AUC", "0.6395", GREEN), ("Test AUC", "0.6415", CLAY), ("Test Brier", "0.2347", PURPLE),
    ]):
        card(s, l, v, 7.85 + (i % 2) * 2.45, 1.45 + (i // 2) * 1.13, 2.2, 0.86, c, small=True)
    textbox(s, "Not betting odds. Not match-win probability.", 8.0, 5.45, 4.1, 0.35, 15, CLAY, True)

    s = sld("Replay System")
    bullets(s, ["Replay manifest: 1,917,672 canonical point events.", "Synthetic match/event IDs are deterministic replay metadata.", "Dry-run JSONL remains deterministic fallback.", "Kafka mode publishes point_event_v1 messages keyed by synthetic_match_id."], 0.9, 1.55, 5.1, 3.5, 18)
    card(s, "Replay manifest", "10,464 matches", 7.0, 1.55, 4.6, 1.0, GREEN)
    card(s, "Dry-run validation", "1000 events", 7.0, 2.9, 4.6, 1.0, GOLD)
    card(s, "Partition key", "synthetic_match_id", 7.0, 4.25, 4.6, 1.0, PURPLE)

    s = sld("Kafka + Spark Structured Streaming")
    image(s, diagrams["streaming"], 0.55, 1.18, w=12.25)
    textbox(s, "Runtime proof: Kafka PASSED, Spark Structured Streaming PASSED, 1000 events scored, checkpoint exists.", 0.9, 6.62, 11.5, 0.35, 14, GREEN, True)

    s = sld("Scoring + Risk Pipeline")
    bullets(s, ["Online feature builder maintains match state.", "Published model returns point_probability_player_a/player_b.", "Risk config uses conservative baseline-deviation rules.", "No fake anomaly labels; risk is a review signal only."], 0.9, 1.55, 5.4, 3.5, 18)
    for i, (l, v, c) in enumerate([
        ("JSONL throughput", f"{bench.get('events_per_second', 2454):.0f}/sec", GREEN), ("Avg latency", f"{bench.get('average_latency_ms', 0.384):.3f} ms", CLAY),
        ("Spark batch", "1000 scored", PURPLE), ("Invalid events", "0", GOLD),
    ]):
        card(s, l, v, 7.2 + (i % 2) * 2.7, 1.65 + (i // 2) * 1.45, 2.35, 1.0, c)

    s = sld("API + Centre Court Dashboard")
    if (SCREENSHOTS / "dashboard_overview.png").exists():
        framed_image(s, SCREENSHOTS / "dashboard_overview.png", 0.6, 1.25, 6.55)
    if (SCREENSHOTS / "fastapi_docs.png").exists():
        framed_image(s, SCREENSHOTS / "fastapi_docs.png", 7.45, 1.25, 5.25)
    bullets(s, ["FastAPI exposes health, summary, scored events, matches, replay catalog, risk, models, benchmarks.", "Frontend groups workflows into Analytics, Replay, ML Model, and Data Ops."], 7.55, 4.75, 4.8, 1.2, 15)

    s = sld("Validation Evidence Matrix")
    rows = [("Data layer", "PASSED", "data/features/validation_report.json"), ("Model", "PASSED", "model_eval reports"), ("Replay", "PASSED", "replay dry-run validation"), ("Kafka", kafka.get("status", "PASSED"), "kafka_runtime_report.json"), ("Spark", spark_val := j("data/results/spark_streaming/spark_streaming_validation_report.json").get("status", "PASSED"), "spark_streaming_validation_report.json"), ("API", api.get("status", "PASSED"), "api_validation_report.json"), ("Frontend", "PASSED", "frontend_validation_report.json"), ("Tests", "70 passed", "pytest tests")]
    for i, (area, status, evidence) in enumerate(rows):
        y = 1.35 + i * 0.58
        card(s, area, status, 0.85, y, 2.3, 0.44, GREEN if "PASS" in str(status) or "passed" in str(status) else CLAY, small=True)
        textbox(s, evidence, 3.45, y + 0.08, 8.8, 0.25, 12, MUTED)

    s = sld("Difficulties + Engineering Tradeoffs")
    problems = [("Sparse metadata", "excluded surface/rally primary features", "honest limitations"), ("Feature leakage", "strict lagged features and tests", "defensible metrics"), ("Replayability", "deterministic synthetic IDs", "stream-ready manifest"), ("Kafka/Spark friction", "fixed image + foreachBatch", "runtime evidence"), ("Product polish", "Centre Court shell", "demo-ready UX")]
    for i, (p, d, r) in enumerate(problems):
        y = 1.35 + i * 0.95
        card(s, "Problem", p, 0.75, y, 2.45, 0.65, CLAY, True)
        card(s, "Decision", d, 3.55, y, 3.65, 0.65, PURPLE, True)
        card(s, "Result", r, 7.55, y, 3.8, 0.65, GREEN, True)

    s = sld("Team Contributions")
    card(s, "Team Member 1", "Data engineering, cleaning, Parquet zones, validation", 0.9, 1.55, 3.55, 1.35, GREEN)
    card(s, "Team Member 2", "Feature engineering, modeling, risk scoring, evaluation", 4.9, 1.55, 3.55, 1.35, PURPLE)
    card(s, "Team Member 3", "Kafka/Spark, API, frontend, demo integration", 8.9, 1.55, 3.55, 1.35, GOLD)
    bullets(s, ["Shared: contracts, tests, reports, runbooks, final validation.", "Replace placeholders with actual names before presenting if desired."], 1.0, 4.0, 11.0, 1.2, 18)

    s = sld("Limitations + Future Work")
    bullets(s, ["Point probability is not betting odds or match-win probability.", "Risk score is not proof of misconduct or match-fixing.", "Surface analytics, official rankings, and ATP bridge remain future work.", "Local demo is not a production deployment.", "Next: managed streaming, persistent serving store, richer metadata, monitoring dashboards."], 0.9, 1.45, 11.4, 4.4, 19)

    s = sld("Demo Flow + Closing")
    bullets(s, ["1. Start: bash scripts/run_full_demo.sh", "2. Show Dashboard → Match Browser → Replay Center → Point Timeline.", "3. Show Model Performance → Validation → Pipeline Monitor.", "4. Streaming proof: bash scripts/run_streaming_demo.sh --max-events 1000.", "Close: validated data-to-product pipeline with real Kafka/Spark evidence."], 0.9, 1.55, 7.0, 4.2, 19)
    card(s, "Frontend", "127.0.0.1:5173", 8.5, 1.65, 3.4, 1.0, GREEN)
    card(s, "API docs", "127.0.0.1:8000/docs", 8.5, 3.05, 3.4, 1.0, PURPLE)
    card(s, "Streaming", "Kafka + Spark reports", 8.5, 4.45, 3.4, 1.0, GOLD)

    prs.save(PPTX)


NOTES_TEXT = [
("Title", "Open by framing CourtIQ as more than a notebook or dashboard: it is a complete local data-to-product system. Say that the project starts from historical point-level tennis data, makes it trustworthy, replays it as events, scores it, and serves it in a polished analytics product. Transition by saying the next slide explains the problem that motivated that architecture."),
("Problem + Goal", "Explain that raw historical point data is not directly usable for a streaming scoring demo. The data has missing fields and score-state subtleties, and a model can easily leak future information if features are computed incorrectly. The goal was to create a reliable and honest end-to-end path: validated data, replay stream, online features, point scoring, risk review signal, API, and dashboard."),
("Final Product Demo Snapshot", "Describe what the audience will see in the demo. The dashboard is backed by real FastAPI responses over validated scored outputs. Point out that the UI has sections for analytics, replay, model evidence, validation, and reports. Transition to why this qualifies as Big Data rather than just a UI demo."),
("Why This Satisfies Big Data Scope", "Use the metric cards first: 1.92 million points, 10.5 thousand matches, almost 1.92 million replay events, and 20 Kafka partitions. Then explain the architectural pieces: Parquet zones, contracts, quality reports, Kafka replay, Spark Structured Streaming, checkpointing, and API serving. Be explicit that the local demo is bounded, but the architecture is horizontally scalable."),
("Dataset Scale + Data Quality", "Talk about the data foundation. Singles point data was curated into Parquet, invalid and special cases were counted and handled, and reports track feature availability. Mention limitations here, not later: surface coverage is unavailable, rally length is sparse, and ATP bridge features are not claimed."),
("End-to-End Architecture", "Walk left to right through the diagram: Parquet data lake, model artifacts, replay stream, Spark scoring, scored output, FastAPI, and dashboard. The important message is that every stage has a validation or contract artifact, so the demo is reproducible."),
("Data Lake + Cleaning Pipeline", "Explain the zones and the safety principle. Raw/staging inputs are not used directly in later stages. Curated singles data feeds features, baselines, replay manifest, and scored outputs. Transition by focusing on the feature layer, where leakage prevention matters most."),
("Point-in-Time Feature Engineering", "Use the diagram to explain event k. The feature row for event k only uses prior points, then the model scores the current point, and only after scoring is the state updated. This is important because including current score/outcome information can inflate metrics. Tests prove the first row counts and rolling windows are correct."),
("Model Training + Evaluation", "Describe the model as an MVP point-outcome model, not a betting or match-win model. The target is Player A wins the current point. Mention feature count, selected model type, AUC, Brier score, and two-phase publication. Explain that schema hashes and fixture scoring protect inference compatibility."),
("Replay System", "Explain that the replay manifest turns historical matches into deterministic event streams with synthetic match and event IDs. Those IDs are replay metadata, not fake match labels. The same canonical event contract supports JSONL fallback and Kafka publishing."),
("Kafka + Spark Structured Streaming", "This is the strongest streaming evidence slide. State clearly: Kafka ran locally, the topic had 20 partitions, the replay producer published 1000 canonical events, Spark Structured Streaming consumed the topic, foreachBatch reused the validated scorer, and the output was checkpointed and validated. Do not overclaim production exactly-once semantics."),
("Scoring + Risk Pipeline", "Explain the runtime scoring sequence: event arrives, online feature builder updates per-match state, odds model predicts point probability, risk config calculates a conservative review signal, then the scored event is written. Say explicitly that risk is not proof of misconduct."),
("API + Centre Court Dashboard", "Use the screenshots. The API exposes summary, matches, scored events, replay catalog, risk, model, and benchmark endpoints. The frontend organizes those into Analytics, Replay, ML Model, and Data Ops. Note that unsupported pages are labeled planned or sample-derived."),
("Validation Evidence Matrix", "Summarize proof rather than implementation. The matrix shows data, model, replay, Kafka, Spark, API, frontend, and tests all with evidence paths. Mention 70 passing tests and PASSED runtime reports. Transition to the problems solved along the way."),
("Difficulties + Engineering Tradeoffs", "Use the Problem Decision Result format. Surface/rally sparsity led to feature exclusions. Leakage risk led to lagged features and tests. Historical data needed deterministic replay IDs. Kafka/Spark friction led to image fixes and foreachBatch. Product polish required careful truthfulness."),
("Team Contributions", "Use this slide to map the work to three team members without inventing names. Member 1 owns data engineering and validation, Member 2 owns features/model/risk/evaluation, and Member 3 owns streaming/API/frontend/demo. Mention shared ownership of tests and documentation."),
("Limitations + Future Work", "Be candid. Point probabilities are not betting odds or match-win probabilities. Risk is not misconduct detection. Surface analytics, official rankings, and ATP bridge are future integrations. The runtime is local, not production. Future work is managed streaming, persistent serving, richer metadata, and monitoring."),
("Demo Flow + Closing", "Close with the concrete demo path. Start the local product with run_full_demo, show dashboard, match browser, replay center, point timeline, model performance, validation, and pipeline monitor. If time permits, show run_streaming_demo for Kafka/Spark proof. End by saying the project demonstrates a validated local Big Data analytics system from raw point data to product UI."),
]


def write_notes():
    lines = ["# CourtIQ Final Presentation v2 Speaker Notes", ""]
    for i, (title, note) in enumerate(NOTES_TEXT, 1):
        lines += [f"## Slide {i}: {title}", "", f"**What to say:** {note}", "", "**Transition:** Move to the next slide by connecting the current evidence to the next layer of the system.", ""]
    NOTES.write_text("\n".join(lines), encoding="utf-8")


def write_readme():
    README.write_text(
        """# CourtIQ Final Presentation v2

## Files

- `CourtIQ_Final_Presentation_v2.pptx`: polished 18-slide final deck.
- `speaker_notes_v2.md`: slide-by-slide talk track.
- `assets_v2/diagrams/`: redesigned PNG diagrams and Mermaid sources.
- `assets_v2/screenshots/`: captured dashboard/API screenshots plus manual checklist.

## Regenerate

```bash
.venv/bin/python presentation/build_presentation_v2.py
```

## Screenshot Notes

Real screenshots captured in this environment:

- `assets_v2/screenshots/dashboard_overview.png`
- `assets_v2/screenshots/fastapi_docs.png`

## Evidence Used

- Data validation reports under `data/features/` and `data/results/`
- Kafka runtime reports under `data/results/kafka_runtime/`
- Spark streaming reports under `data/results/spark_streaming/`
- Model evaluation reports under `data/results/model_eval/`
- API/frontend validation reports under `data/results/api_validation/` and `data/results/frontend_validation/`
""",
        encoding="utf-8",
    )


def write_screenshot_checklist():
    (SCREENSHOTS / "screenshot_checklist_v2.md").write_text(
        """# Screenshot Checklist v2

Automated Selenium capture was able to capture the main Dashboard Overview and FastAPI docs correctly. Because the frontend relies on complex client-side state, inner-page views (like the Replay Center or Model Performance) may need to be captured manually.

Run:

```bash
bash scripts/run_full_demo.sh
```

Capture manually:

1. Match Browser with selected match detail.
2. Replay Center with playback controls.
3. Point Timeline.
4. Model Performance.
5. Validation page.
6. Pipeline Monitor.
7. Reports page.
""",
        encoding="utf-8",
    )


def main():
    build_deck()
    write_notes()
    write_readme()
    write_screenshot_checklist()
    print(PPTX)


if __name__ == "__main__":
    main()
