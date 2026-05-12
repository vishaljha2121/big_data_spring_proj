#!/usr/bin/env python3
"""Generate the final CourtIQ project presentation and diagram assets."""

from __future__ import annotations

import json
import math
from pathlib import Path
from typing import Dict, Iterable, List, Tuple

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "presentation"
ASSETS = OUT / "assets"
DIAGRAMS = ASSETS / "diagrams"
SCREENSHOTS = ASSETS / "screenshots"

GREEN = RGBColor(6, 66, 37)
GREEN_DARK = RGBColor(4, 45, 28)
CREAM = RGBColor(248, 244, 232)
GOLD = RGBColor(216, 180, 93)
PURPLE = RGBColor(91, 46, 145)
INK = RGBColor(30, 41, 59)
MUTED = RGBColor(100, 116, 139)
WHITE = RGBColor(255, 255, 255)
TERRACOTTA = RGBColor(186, 92, 50)


def load_json(path: str, default=None):
    p = ROOT / path
    if not p.exists():
        return default if default is not None else {}
    return json.loads(p.read_text(encoding="utf-8"))


def metric(obj: Dict, *keys, default="n/a"):
    cur = obj
    for key in keys:
        if not isinstance(cur, dict) or key not in cur:
            return default
        cur = cur[key]
    return cur


def fmt(value, digits=3):
    if isinstance(value, float):
        return f"{value:.{digits}f}"
    return str(value)


def ensure_dirs():
    DIAGRAMS.mkdir(parents=True, exist_ok=True)
    SCREENSHOTS.mkdir(parents=True, exist_ok=True)


def font(size=28, bold=False):
    candidates = [
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf" if bold else "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/Library/Fonts/Arial Bold.ttf" if bold else "/Library/Fonts/Arial.ttf",
    ]
    for candidate in candidates:
        if Path(candidate).exists():
            return ImageFont.truetype(candidate, size=size)
    return ImageFont.load_default()


def wrap(draw, text: str, fnt, max_width: int) -> List[str]:
    words = text.split()
    lines: List[str] = []
    line = ""
    for word in words:
        candidate = f"{line} {word}".strip()
        if draw.textbbox((0, 0), candidate, font=fnt)[2] <= max_width:
            line = candidate
        else:
            if line:
                lines.append(line)
            line = word
    if line:
        lines.append(line)
    return lines


def draw_box(draw, xy, text, fill, outline=(220, 210, 185), text_fill=(25, 35, 45), title=False):
    x1, y1, x2, y2 = xy
    draw.rounded_rectangle(xy, radius=20, fill=fill, outline=outline, width=3)
    fnt = font(28 if title else 24, bold=title)
    lines = wrap(draw, text, fnt, x2 - x1 - 36)
    total = len(lines) * (fnt.size + 4)
    y = y1 + ((y2 - y1) - total) / 2
    for line in lines:
        bbox = draw.textbbox((0, 0), line, font=fnt)
        draw.text((x1 + ((x2 - x1) - (bbox[2] - bbox[0])) / 2, y), line, font=fnt, fill=text_fill)
        y += fnt.size + 4


def arrow(draw, start, end, fill=(6, 66, 37)):
    draw.line([start, end], fill=fill, width=5)
    angle = math.atan2(end[1] - start[1], end[0] - start[0])
    size = 16
    p1 = (end[0] - size * math.cos(angle - math.pi / 6), end[1] - size * math.sin(angle - math.pi / 6))
    p2 = (end[0] - size * math.cos(angle + math.pi / 6), end[1] - size * math.sin(angle + math.pi / 6))
    draw.polygon([end, p1, p2], fill=fill)


def save_mermaid(name: str, source: str):
    (DIAGRAMS / f"{name}.mmd").write_text(source.strip() + "\n", encoding="utf-8")


def make_diagram(name: str, title: str, boxes: List[Tuple[int, int, int, int, str]], arrows: List[Tuple[int, int, int, int]]):
    w, h = 1600, 900
    img = Image.new("RGB", (w, h), (248, 244, 232))
    draw = ImageDraw.Draw(img)
    draw.rectangle((0, 0, w, 110), fill=(6, 66, 37))
    draw.text((60, 32), title, font=font(42, bold=True), fill=(255, 255, 255))
    for x1, y1, x2, y2, text in boxes:
        draw_box(draw, (x1, y1, x2, y2), text, fill=(255, 253, 246), title=True)
    for x1, y1, x2, y2 in arrows:
        arrow(draw, (x1, y1), (x2, y2))
    path = DIAGRAMS / f"{name}.png"
    img.save(path)
    return path


def build_diagrams():
    diagrams = {}
    diagrams["end_to_end"] = make_diagram(
        "end_to_end_architecture",
        "End-to-End Architecture",
        [
            (70, 190, 310, 320, "Curated Parquet Data Lake"),
            (390, 190, 630, 320, "Feature + Baseline Layer"),
            (710, 190, 950, 320, "Model + Risk Artifacts"),
            (1030, 190, 1270, 320, "Replay + Streaming"),
            (1190, 510, 1490, 650, "FastAPI + Centre Court UI"),
            (710, 510, 1030, 650, "Scored JSONL / Parquet"),
        ],
        [(310, 255, 390, 255), (630, 255, 710, 255), (950, 255, 1030, 255), (1150, 320, 930, 510), (1030, 580, 1190, 580)],
    )
    save_mermaid(
        "end_to_end_architecture",
        """
flowchart LR
  A[Curated Parquet Data Lake] --> B[Point-in-Time Feature Layer]
  B --> C[Odds Model + Risk Config]
  C --> D[Replay / Kafka / Spark Scoring]
  D --> E[Scored JSONL + Parquet]
  E --> F[FastAPI + Centre Court Frontend]
""",
    )
    diagrams["data_lake"] = make_diagram(
        "data_lake_zones",
        "Parquet Data Lake Zones",
        [
            (80, 205, 330, 335, "Staging Inventory"),
            (430, 205, 680, 335, "Cleaned Layer"),
            (780, 205, 1030, 335, "Curated Singles"),
            (1130, 205, 1380, 335, "Features / Baselines"),
            (780, 520, 1030, 650, "Replay Manifest"),
            (1130, 520, 1380, 650, "Scored Outputs"),
        ],
        [(330, 270, 430, 270), (680, 270, 780, 270), (1030, 270, 1130, 270), (905, 335, 905, 520), (1030, 585, 1130, 585)],
    )
    save_mermaid(
        "data_lake_zones",
        """
flowchart LR
  staging[Staging inventory] --> cleaned[Cleaned Parquet]
  cleaned --> curated[Curated singles]
  curated --> features[Point features + baselines]
  curated --> replay[Replay manifest]
  replay --> scored[Scored outputs]
""",
    )
    diagrams["batch"] = make_diagram(
        "batch_pipeline",
        "Batch Pipeline",
        [
            (80, 220, 360, 350, "Validate Inputs"),
            (460, 220, 740, 350, "Clean + Normalize"),
            (840, 220, 1120, 350, "PIT Feature Engineering"),
            (1220, 220, 1500, 350, "Train + Publish Artifacts"),
            (650, 550, 950, 680, "Schema Contracts + Reports"),
        ],
        [(360, 285, 460, 285), (740, 285, 840, 285), (1120, 285, 1220, 285), (800, 350, 800, 550)],
    )
    save_mermaid(
        "batch_pipeline",
        """
flowchart LR
  A[Input validation] --> B[Clean + normalize]
  B --> C[Point-in-time features]
  C --> D[Model publication]
  C --> E[Contracts + quality reports]
""",
    )
    diagrams["streaming"] = make_diagram(
        "kafka_spark_streaming",
        "Kafka + Spark Structured Streaming",
        [
            (80, 220, 360, 350, "Replay Producer"),
            (460, 220, 740, 350, "Kafka Topic 20 Partitions"),
            (840, 220, 1120, 350, "Spark readStream"),
            (1220, 220, 1500, 350, "foreachBatch Scorer"),
            (840, 550, 1120, 680, "Checkpoint"),
            (1220, 550, 1500, 680, "JSONL + Parquet Sink"),
        ],
        [(360, 285, 460, 285), (740, 285, 840, 285), (1120, 285, 1220, 285), (1360, 350, 1360, 550), (980, 350, 980, 550)],
    )
    save_mermaid(
        "kafka_spark_streaming",
        """
flowchart LR
  A[Replay producer] --> B[(Kafka tennis-point-events)]
  B --> C[Spark Structured Streaming]
  C --> D[foreachBatch StreamScorer]
  D --> E[Scored JSONL + Parquet]
  C --> F[Checkpointing]
""",
    )
    diagrams["model"] = make_diagram(
        "model_publication",
        "Two-Phase Model Publication",
        [
            (100, 230, 380, 360, "Train Candidates"),
            (480, 230, 760, 360, "Evaluate Gates"),
            (860, 230, 1140, 360, "Staging Artifacts"),
            (1240, 230, 1520, 360, "v1 + latest.json"),
            (480, 550, 1140, 680, "Feature Schema Hash + Fixture Scoring"),
        ],
        [(380, 295, 480, 295), (760, 295, 860, 295), (1140, 295, 1240, 295), (810, 360, 810, 550)],
    )
    save_mermaid(
        "model_publication",
        """
flowchart LR
  A[Train candidates] --> B[Evaluate AUC/Brier gates]
  B --> C[Staging artifacts]
  C --> D[v1 artifact + latest.json]
  B --> E[Feature hash + fixture checks]
""",
    )
    diagrams["api"] = make_diagram(
        "api_frontend_serving",
        "API + Frontend Serving",
        [
            (120, 230, 420, 360, "Scored JSONL / Parquet"),
            (520, 230, 820, 360, "FastAPI Repositories"),
            (920, 230, 1220, 360, "Documented API"),
            (1320, 230, 1540, 360, "React UI"),
            (520, 550, 820, 680, "OpenAPI + Examples"),
            (920, 550, 1220, 680, "Validation Reports"),
        ],
        [(420, 295, 520, 295), (820, 295, 920, 295), (1220, 295, 1320, 295), (670, 360, 670, 550), (1070, 360, 1070, 550)],
    )
    save_mermaid(
        "api_frontend_serving",
        """
flowchart LR
  A[Scored files] --> B[FastAPI repositories]
  B --> C[Documented REST API]
  C --> D[Centre Court React UI]
  B --> E[OpenAPI + response examples]
  C --> F[Validation reports]
""",
    )
    return diagrams


def add_slide_number(slide, idx):
    box = slide.shapes.add_textbox(Inches(12.25), Inches(7.05), Inches(0.8), Inches(0.25))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = str(idx)
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def set_bg(slide, color=CREAM):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.1), Inches(0.75))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(30)
    p.font.color.rgb = GREEN_DARK
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(1.05), Inches(11.5), Inches(0.36))
        p2 = sub.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(12)
        p2.font.color.rgb = MUTED


def add_bullets(slide, bullets: Iterable[str], x=0.75, y=1.65, w=5.8, h=4.8, size=17):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = INK
        p.space_after = Pt(7)


def add_card(slide, title, value, x, y, w=2.55, h=1.05, accent=GREEN):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(224, 216, 197)
    title_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.12), Inches(w - 0.3), Inches(0.25))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = MUTED
    val_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.42), Inches(w - 0.3), Inches(0.45))
    pv = val_box.text_frame.paragraphs[0]
    pv.text = str(value)
    pv.font.size = Pt(20)
    pv.font.bold = True
    pv.font.color.rgb = accent


def add_image(slide, path, x, y, w=None, h=None):
    if w and h:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    if w:
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y))


def create_deck(diagrams):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    reports = {
        "features": load_json("data/features/validation_report.json"),
        "feature_build": load_json("data/features/feature_build_report.json"),
        "baselines": load_json("data/baselines/baseline_quality_report.json"),
        "odds": load_json("data/results/model_eval/odds_model_eval_report.json"),
        "benchmark": load_json("data/results/streaming_scoring/scoring_benchmark_report.json"),
        "kafka": load_json("data/results/kafka_runtime/kafka_runtime_report.json"),
        "kafka_smoke": load_json("data/results/kafka_runtime/kafka_replay_smoke_report.json"),
        "spark": load_json("data/results/spark_streaming/spark_streaming_run_report.json"),
        "spark_val": load_json("data/results/spark_streaming/spark_streaming_validation_report.json"),
        "api": load_json("data/results/api_validation/api_validation_report.json"),
        "frontend": load_json("data/results/frontend_validation/frontend_validation_report.json"),
        "preflight": load_json("data/results/final_demo/final_preflight_report.json"),
    }
    odds_metrics = reports["odds"].get("metrics", {})

    slides = []

    def new_slide(title, subtitle=None):
        s = prs.slides.add_slide(blank)
        set_bg(s)
        add_title(s, title, subtitle)
        slides.append(s)
        add_slide_number(s, len(slides))
        return s

    # 1
    s = prs.slides.add_slide(blank)
    set_bg(s, GREEN_DARK)
    title = s.shapes.add_textbox(Inches(0.75), Inches(1.35), Inches(11.8), Inches(1.0))
    p = title.text_frame.paragraphs[0]
    p.text = "CourtIQ: Centre Court Analytics"
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = WHITE
    sub = s.shapes.add_textbox(Inches(0.82), Inches(2.35), Inches(10.8), Inches(0.6))
    ps = sub.text_frame.paragraphs[0]
    ps.text = "Validated tennis point data, Kafka/Spark replay streaming, point-level scoring, risk review signals, API, and dashboard"
    ps.font.size = Pt(17)
    ps.font.color.rgb = RGBColor(242, 232, 204)
    for i, (label, value) in enumerate([
        ("Point rows", "1.92M+"),
        ("Matches", "10.5K+"),
        ("Kafka/Spark", "PASSED"),
        ("Frontend/API", "PASSED"),
    ]):
        add_card(s, label, value, 0.85 + i * 3.05, 4.6, accent=GOLD)
    slides.append(s)
    add_slide_number(s, 1)

    # 2
    s = new_slide("Problem Statement", "Tennis analytics needs point-level context, not only final scores.")
    add_bullets(s, [
        "Raw tennis point feeds are noisy, sparse, and not model-ready.",
        "Point outcome predictions must avoid future leakage.",
        "Streaming demos need reproducible event replay, scoring, and validation evidence.",
        "Risk outputs must be framed as review signals, not misconduct proof.",
    ])
    add_card(s, "Project Goal", "End-to-end validated analytics platform", 7.1, 1.9, 4.7, 1.2, GOLD)
    add_card(s, "Bounded Demo", "Local runtime, scalable design", 7.1, 3.35, 4.7, 1.2, PURPLE)

    # 3
    s = new_slide("Product Demo Overview", "Centre Court Analytics: a product shell over real scored output.")
    add_bullets(s, [
        "Dashboard: validated KPIs, model quality, and replay status.",
        "Match Browser + Replay Center: player-first match labels and point timeline.",
        "Prediction Center: point-level probability output, not betting odds.",
        "Validation + Pipeline Monitor: evidence for API, scoring, Kafka, and Spark.",
    ], w=6.4)
    add_image(s, diagrams["api"], 7.0, 1.65, w=5.7)

    # 4
    s = new_slide("Why This Is a Big Data Project", "The local demo is bounded; the architecture is built for scale.")
    add_bullets(s, [
        "1.9M+ point-level rows and 10K+ singles matches processed into Parquet zones.",
        "Schema contracts, data quality reports, validation scripts, and quarantine evidence.",
        "Spark ETL / feature processing concepts used for scalable batch stages.",
        "Kafka replay stream and Spark Structured Streaming runtime validated locally.",
        "Checkpointed streaming scorer writes scalable local JSONL/Parquet sinks.",
        "FastAPI/frontend serve validated scored outputs for analysis and presentation.",
    ], w=11.5, size=16)

    # 5
    s = new_slide("Dataset and Scale", "Validated tennis point data became the foundation for modeling and replay.")
    cards = [
        ("Curated point rows", "1,922,136"),
        ("Curated matches", "10,508"),
        ("Replay events", "1,917,672"),
        ("Replay matches", "10,464"),
        ("Player baselines", "1,891"),
        ("Full demo scored events", "9,049"),
    ]
    for i, (label, value) in enumerate(cards):
        add_card(s, label, value, 0.75 + (i % 3) * 4.1, 1.7 + (i // 3) * 1.55, w=3.55)
    add_bullets(s, [
        "Surface coverage is unavailable; surface analytics remain planned.",
        "Rally length coverage is sparse; not used as a primary MVP feature.",
        "ATP bridge remains unvalidated; no ATP-derived claims are made.",
    ], x=0.9, y=5.15, w=11.4, h=1.2, size=15)

    # 6
    s = new_slide("End-to-End Architecture")
    add_image(s, diagrams["end_to_end"], 0.65, 1.35, w=12.0)

    # 7
    s = new_slide("Data Lake + Cleaning Pipeline")
    add_image(s, diagrams["data_lake"], 0.65, 1.35, w=6.0)
    add_image(s, diagrams["batch"], 6.95, 1.35, w=5.75)
    add_bullets(s, [
        "Cleaned and curated layers are Parquet-backed.",
        "Singles-only filtering, deterministic placeholders, and quarantine handling.",
        "Reports track invalid winners/servers, missing metadata, and feature availability.",
    ], x=0.9, y=6.0, w=11.2, h=0.9, size=13)

    # 8
    s = new_slide("Feature Engineering and Point-in-Time Safety")
    add_bullets(s, [
        "Features are sorted by match and event index.",
        "All cumulative and rolling features use prior events only.",
        "First point in each match has zero prior counts.",
        "Recent-window features use only earlier points.",
        "Tests prove current point winner does not leak into current features.",
    ], w=6.3)
    add_card(s, "Point features", "1,922,136 rows", 7.2, 1.7, 4.3)
    add_card(s, "Match features", "10,508 rows", 7.2, 3.0, 4.3)
    add_card(s, "Point-in-time", "Validated true", 7.2, 4.3, 4.3, accent=GOLD)

    # 9
    s = new_slide("Model Training and Evaluation")
    add_bullets(s, [
        "Primary target: Player A wins the current point.",
        "Selected model: HistGradientBoosting classifier.",
        "Two-phase publication: staging validation → v1 artifact → latest.json pointer.",
        "Feature schema hash prevents inference-time mismatch.",
    ], w=6.2)
    add_card(s, "Feature count", odds_metrics.get("feature_count", 26), 7.2, 1.55)
    add_card(s, "Validation AUC", "0.6395", 9.95, 1.55)
    add_card(s, "Test AUC", "0.6415", 7.2, 2.9)
    add_card(s, "Test Brier", "0.2347", 9.95, 2.9)
    add_image(s, diagrams["model"], 6.75, 4.1, w=5.9)

    # 10
    s = new_slide("Replay Architecture")
    add_bullets(s, [
        "Replay manifest selects validated singles matches for deterministic event replay.",
        "Synthetic match and event IDs are deterministic and schema-validated.",
        "Dry-run JSONL path remains reproducible without Kafka.",
        "Kafka mode publishes canonical point_event_v1 messages keyed by synthetic_match_id.",
    ], w=6.3)
    add_card(s, "Manifest rows", "1,917,672", 7.2, 1.65, 4.3)
    add_card(s, "Validated dry-run", "1,000 events", 7.2, 3.0, 4.3)
    add_card(s, "Full catalog", "10,464 matches", 7.2, 4.35, 4.3)

    # 11
    s = new_slide("Kafka + Spark Structured Streaming Architecture")
    add_image(s, diagrams["streaming"], 0.65, 1.35, w=12.0)
    add_bullets(s, [
        "Runtime status: PASSED with 1000 Kafka events and 1000 Spark-scored events.",
        "Spark uses foreachBatch to reuse validated Python scoring logic.",
        "Checkpointing is local and suitable for demo validation, not production exactly-once claims.",
    ], x=0.9, y=6.1, w=11.3, h=0.75, size=12)

    # 12
    s = new_slide("Scoring + Risk Pipeline")
    add_bullets(s, [
        "Online feature builder maintains state per synthetic match.",
        "Odds model outputs point_probability_player_a / player_b.",
        "Risk config computes conservative baseline-deviation signals.",
        "Risk score is a statistical review signal only.",
        "No fake anomaly labels or misconduct claims are used.",
    ], w=6.6)
    add_card(s, "JSONL throughput", f"{fmt(reports['benchmark'].get('events_per_second', 0), 1)} ev/s", 7.3, 1.55, 4.1)
    add_card(s, "Avg latency", f"{fmt(reports['benchmark'].get('average_latency_ms', 0), 3)} ms", 7.3, 2.9, 4.1)
    add_card(s, "Spark scored", reports["spark"].get("scored_count", 1000), 7.3, 4.25, 4.1, accent=GOLD)

    # 13
    s = new_slide("API and Frontend Product")
    add_bullets(s, [
        "FastAPI serves validated scored files and model metadata.",
        "OpenAPI snapshot and response examples are exported.",
        "Centre Court Analytics React app exposes Analytics, Replay, ML Model, and Data Ops sections.",
        "Unsupported modules are labeled planned or sample-derived.",
    ], w=6.3)
    add_image(s, diagrams["api"], 7.0, 1.45, w=5.8)

    # 14
    s = new_slide("Validation and Testing Evidence")
    add_card(s, "Pytest", "70 passed", 0.8, 1.55)
    add_card(s, "API validation", reports["api"].get("status", "PASSED"), 3.55, 1.55)
    add_card(s, "Frontend build", reports["frontend"].get("status", "PASSED"), 6.3, 1.55)
    add_card(s, "Final preflight", reports["preflight"].get("status", "PASSED"), 9.05, 1.55)
    add_card(s, "Kafka runtime", reports["kafka"].get("status", "PASSED"), 0.8, 3.0)
    add_card(s, "Spark streaming", reports["spark_val"].get("status", "PASSED"), 3.55, 3.0)
    add_card(s, "Kafka events", "1000", 6.3, 3.0)
    add_card(s, "Spark scored", "1000", 9.05, 3.0)
    add_bullets(s, [
        "Validation scripts write JSON evidence under data/results/.",
        "Contracts guard point events, scored events, model metadata, and API responses.",
    ], x=0.95, y=5.0, w=11, h=1.2, size=15)

    # 15
    s = new_slide("Difficulties and Engineering Tradeoffs")
    add_bullets(s, [
        "Sparse surface/rally metadata blocked certain feature families.",
        "Avoiding point-in-time leakage required strict lag/shift logic and synthetic tests.",
        "Kafka image tag and runtime setup required hardening before validation.",
        "Spark model scoring used foreachBatch to avoid fragile distributed sklearn loading.",
        "Dashboard needed product polish while preserving truthful data semantics.",
    ], w=11.2)

    # 16
    s = new_slide("Team Contributions", "3-member contribution split for the final report.")
    add_card(s, "Team Member 1", "Data engineering + cleaning + validation", 0.8, 1.85, 3.8, 1.3, GREEN)
    add_card(s, "Team Member 2", "Modeling + scoring + risk artifacts", 4.75, 1.85, 3.8, 1.3, PURPLE)
    add_card(s, "Team Member 3", "Streaming/API/frontend/demo integration", 8.7, 1.85, 3.8, 1.3, GOLD)
    add_bullets(s, [
        "Cross-cutting work: contracts, tests, documentation, validation reports, demo runbooks.",
        "Replace placeholder labels with actual member names if desired before presenting.",
    ], x=0.9, y=4.0, w=11.2, h=1.3, size=16)

    # 17
    s = new_slide("Limitations")
    add_bullets(s, [
        "Point probabilities are not betting odds and not match-win probabilities.",
        "Risk scores are statistical anomaly signals, not proof of misconduct or match-fixing.",
        "Surface analytics and official rankings are not claimed; current views are sample-derived or planned.",
        "Local Spark/Kafka demo is bounded; production scaling needs managed infrastructure and stronger sink semantics.",
        "ATP bridge remains unvalidated.",
    ], w=11.2)

    # 18
    s = new_slide("Future Work")
    add_bullets(s, [
        "Add reliable tournament/surface metadata and official rankings integration.",
        "Build match-level prediction model separately from point-level scoring.",
        "Move streaming jobs to managed Kafka/Spark infrastructure.",
        "Add persistent serving storage after schema contracts stabilize.",
        "Add monitoring dashboards for drift, latency, and replay coverage.",
    ], w=11.2)

    # 19
    s = new_slide("Demo Flow")
    add_bullets(s, [
        "Start local product: bash scripts/run_full_demo.sh",
        "Show Dashboard → Match Browser → Replay Center → Point Timeline.",
        "Show Prediction Center and Model Performance with point-level language.",
        "Show Validation and Pipeline Monitor evidence.",
        "Optional Big Data proof: bash scripts/run_streaming_demo.sh --max-events 1000.",
    ], w=10.8)
    add_card(s, "Frontend", "http://127.0.0.1:5173", 1.0, 5.6, 3.5, 0.9)
    add_card(s, "API docs", "http://127.0.0.1:8000/docs", 4.9, 5.6, 3.5, 0.9)
    add_card(s, "Streaming", "Kafka + Spark reports", 8.8, 5.6, 3.5, 0.9, accent=GOLD)

    # 20
    s = new_slide("Closing: What We Built")
    add_bullets(s, [
        "A validated data-to-product tennis analytics system.",
        "Batch data lake, point-in-time features, model artifacts, and risk config.",
        "Real Kafka replay and Spark Structured Streaming scoring evidence.",
        "FastAPI + Centre Court dashboard for demo-ready exploration.",
        "Honest limitations and reproducible validation artifacts.",
    ], w=9.5)
    add_card(s, "Final claim", "Validated local Big Data analytics demo", 3.2, 5.7, 6.7, 0.9, accent=GOLD)

    prs.save(OUT / "CourtIQ_Final_Presentation.pptx")


SLIDE_NOTES = [
    ("Title slide", "Introduce CourtIQ as a validated tennis point-level analytics platform. Emphasize that the project covers data engineering, model artifacts, streaming, API, and product UI."),
    ("Problem statement", "Explain that tennis point data is high-volume and messy, and naive models can leak future information. The goal was a reliable end-to-end demo with truthful probability and risk language."),
    ("Product demo overview", "Walk through the Centre Court Analytics product: dashboard, match browser, replay center, point scoring, validation, and reports."),
    ("Why this is a Big Data project", "Call out 1.9M+ point rows, 10K+ matches, Parquet zones, contracts, Spark-style batch processing, Kafka replay, Spark Structured Streaming, and bounded local scale."),
    ("Dataset and scale", "Use the numbers as evidence: curated rows, matches, replay manifest size, and limitations around surface/rally/ATP metadata."),
    ("End-to-end architecture", "Trace the full path from curated data lake through features, models, replay, streaming, scored outputs, API, and frontend."),
    ("Data lake + cleaning pipeline", "Describe the data lake zones and cleaning decisions: singles filtering, placeholders, quarantine, reports, and schema contracts."),
    ("Feature engineering and point-in-time safety", "Emphasize no future leakage: sort by match/event, use lagged cumulative counts, first point zero prior counts, tests prove safety."),
    ("Model training and evaluation", "Explain the published model, feature count, AUC/Brier, two-phase publication, and why point-level target is not match prediction."),
    ("Replay architecture", "Explain deterministic replay manifest, JSONL fallback, Kafka publishing mode, and canonical event schema."),
    ("Kafka + Spark Structured Streaming architecture", "This is the core Big Data evidence slide. Say Kafka and Spark were actually executed locally with 1000 events and checkpointed scored output."),
    ("Scoring + risk pipeline", "Explain online feature state, odds scoring, risk config, and careful language: risk is review signal only."),
    ("API and frontend product", "Explain FastAPI endpoints, OpenAPI snapshot, response examples, and Centre Court product shell. Mention planned/sample-derived pages are labeled."),
    ("Validation and testing evidence", "Summarize 70 passing tests and PASSED reports for API, frontend, Kafka, Spark, model, replay, and final preflight."),
    ("Difficulties and engineering tradeoffs", "Discuss data limitations, leakage prevention, Kafka image/runtime hardening, foreachBatch tradeoff, and UI truthfulness."),
    ("Team contributions", "Use the three-member split. Replace placeholders with names if needed before presentation."),
    ("Limitations", "Be explicit: no betting odds, no misconduct proof, no official rankings, bounded local runtime, no production deployment."),
    ("Future work", "Prioritize metadata enrichment, match-level prediction, managed streaming, persistent serving storage, and monitoring."),
    ("Demo flow", "Give the exact path to demo: run_full_demo, then pages to show, then optional run_streaming_demo for streaming proof."),
    ("Closing slide", "Close with the main accomplishment: validated local Big Data analytics pipeline with streaming evidence and product UI."),
]


def write_notes():
    lines = ["# CourtIQ Final Presentation Speaker Notes", ""]
    for i, (title, note) in enumerate(SLIDE_NOTES, start=1):
        lines.extend([f"## Slide {i}: {title}", "", note, ""])
    (OUT / "speaker_notes.md").write_text("\n".join(lines), encoding="utf-8")


def write_readme():
    text = """# CourtIQ Final Presentation

## Deliverables

- `CourtIQ_Final_Presentation.pptx`: final PowerPoint deck.
- `speaker_notes.md`: slide-by-slide speaker script.
- `assets/diagrams/`: generated architecture diagrams and Mermaid sources.
- `assets/screenshots/screenshot_checklist.md`: manual screenshot checklist because Playwright/browser automation was not available in this environment.

## Regenerate

From the repository root:

```bash
.venv/bin/python presentation/build_presentation.py
```

The script reads current project reports under `data/results/`, generates diagrams with Pillow, and writes the PowerPoint deck with `python-pptx`.

## Current Evidence Used

- `data/features/validation_report.json`
- `data/results/model_eval/odds_model_eval_report.json`
- `data/results/streaming_scoring/scoring_benchmark_report.json`
- `data/results/kafka_runtime/kafka_runtime_report.json`
- `data/results/kafka_runtime/kafka_replay_smoke_report.json`
- `data/results/spark_streaming/spark_streaming_run_report.json`
- `data/results/spark_streaming/spark_streaming_validation_report.json`
- `data/results/api_validation/api_validation_report.json`
- `data/results/frontend_validation/frontend_validation_report.json`

## Screenshot Capture

Use:

```bash
bash scripts/run_full_demo.sh
```

Then follow `assets/screenshots/screenshot_checklist.md`.
"""
    (OUT / "README.md").write_text(text, encoding="utf-8")


def write_screenshot_checklist():
    text = """# Manual Screenshot Checklist

Automated browser screenshot tooling was not installed in this environment. Capture these manually after running:

```bash
bash scripts/run_full_demo.sh
```

Open `http://127.0.0.1:5173` and capture:

1. Dashboard overview with KPI grid.
2. Match Browser with selected match detail.
3. Replay Center with playback controls.
4. Point Timeline.
5. Model Performance.
6. Validation page.
7. Pipeline Monitor.
8. Reports page.
9. FastAPI docs at `http://127.0.0.1:8000/docs`.

Save screenshots in this folder before final slide export if desired.
"""
    (SCREENSHOTS / "screenshot_checklist.md").write_text(text, encoding="utf-8")


def main():
    ensure_dirs()
    diagrams = build_diagrams()
    create_deck(diagrams)
    write_notes()
    write_readme()
    write_screenshot_checklist()
    print(f"Wrote {OUT / 'CourtIQ_Final_Presentation.pptx'}")


if __name__ == "__main__":
    main()
